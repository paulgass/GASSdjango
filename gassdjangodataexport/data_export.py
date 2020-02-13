from django.http import HttpResponse
import json
import requests
import requests.exceptions
import sys
import os
from io import BytesIO
from ezodf import newdoc, opendoc, Paragraph, Heading, Sheet
from docx import Document
from docx.shared import Inches
import csv
import xlwt
import xlsxwriter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from reportlab.pdfgen import canvas
import random
from PIL import Image, ImageDraw
from gtts import gTTS
import qrcode
from zipfile import ZipFile
from datetime import datetime

array_names = [line.rstrip('\n') for line in open('gassdjangodataexport/names.txt')]
array_nouns = [line.rstrip('\n') for line in open('gassdjangodataexport/nouns.txt')]
array_verbs = [line.rstrip('\n') for line in open('gassdjangodataexport/verbs.txt')]

def get_random_word(array_words):
    return array_words[random.randint(0, (len(array_words) - 1))]

def downloadODT():
    response = HttpResponse(content_type="application/vnd.oasis.opendocument.text")
    response['Content-Disposition'] = 'attachment; filename="ODT_Export.odt"'

    odt = newdoc(doctype='odt', filename='ODT_Export.odt', template=None)
    odt.body += Heading("Chapter 1")
    odt.body += Paragraph("This is a paragraph for gassdjangoproject.")
    odt.save()

    return HttpResponse("/gassdjango-api/gassdjangodataexport/data_export.py downloadODT() FUNCTIONALITY IS NOT YET COMPLETE :(\nHowever, the .odt was generated and is located @ /gassdjango-api/ODT_Export.odt", status=200)
    #return response

def downloadDOC():
    response = HttpResponse(content_type="application/msword")
    response['Content-Disposition'] = 'attachment; filename="DOC_Export.doc"'

    document = Document()

    document.add_heading('Document Title', 0)

    p = document.add_paragraph('A plain paragraph for gassdjangoproject having some ')
    p.add_run('bold').bold = True
    p.add_run(' and some ')
    p.add_run('italic.').italic = True

    document.add_heading('Heading, level 1', level=1)
    document.add_paragraph('Intense quote', style='Intense Quote')

    document.add_paragraph(
        'first item in unordered list', style='List Bullet'
    )
    document.add_paragraph(
        'first item in ordered list', style='List Number'
    )

    #document.add_picture('monty-truth.png', width=Inches(1.25))

    records = (
        (3, '101', 'Spam'),
        (7, '422', 'Eggs'),
        (4, '631', 'Spam, spam, eggs, and spam')
    )

    table = document.add_table(rows=1, cols=3)
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Qty'
    hdr_cells[1].text = 'Id'
    hdr_cells[2].text = 'Desc'
    for qty, id, desc in records:
        row_cells = table.add_row().cells
        row_cells[0].text = str(qty)
        row_cells[1].text = id
        row_cells[2].text = desc

    document.add_page_break()

    document.save(response)
    return response

def downloadDOCX():
    response = HttpResponse(content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    response['Content-Disposition'] = 'attachment; filename="DOCX_Export.docx"'

    document = Document()

    document.add_heading('Document Title', 0)

    p = document.add_paragraph('A plain paragraph for gassdjangoproject having some ')
    p.add_run('bold').bold = True
    p.add_run(' and some ')
    p.add_run('italic.').italic = True

    document.add_heading('Heading, level 1', level=1)
    document.add_paragraph('Intense quote', style='Intense Quote')

    document.add_paragraph(
        'first item in unordered list', style='List Bullet'
    )
    document.add_paragraph(
        'first item in ordered list', style='List Number'
    )

    #document.add_picture('monty-truth.png', width=Inches(1.25))

    records = (
        (3, '101', 'Spam'),
        (7, '422', 'Eggs'),
        (4, '631', 'Spam, spam, eggs, and spam')
    )

    table = document.add_table(rows=1, cols=3)
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Qty'
    hdr_cells[1].text = 'Id'
    hdr_cells[2].text = 'Desc'
    for qty, id, desc in records:
        row_cells = table.add_row().cells
        row_cells[0].text = str(qty)
        row_cells[1].text = id
        row_cells[2].text = desc

    document.add_page_break()

    document.save(response)
    return response

def downloadCSV():
    response = HttpResponse(content_type="text/csv")
    response['Content-Disposition'] = 'attachment; filename="CSV_Export.csv"'
    csv_writer = csv.writer(response)

    csv_data_array = []
    csv_data_array.append(["sunday","monday","tuesday","wednesday","thursday","friday","saturday"])
    csv_data_array.append(["123","345","","","934"])
    csv_data_array.append(["","657","","gassdjangoproject","","092"])
    csv_data_array.append(["","","","234","934",""])

    for a in csv_data_array:
        csv_writer.writerow(a)

    return response
 
def downloadODS():
    response = HttpResponse(content_type='application/vnd.oasis.opendocument.spreadsheet')
    response['Content-Disposition'] = 'attachment; filename="ODS_Export.ods"'

    ods = newdoc(doctype='ods', filename='ODS_Export.ods', template=None)
    sheet = Sheet('gassdjangoproject', size=(10, 10))
    ods.sheets += sheet
    sheet['A1'].set_value("Chapter 2")
    sheet['B2'].set_value(3.141592)
    sheet['C3'].set_value(100, currency='USD')
    sheet['D4'].formula = "of:=SUM([.B2];[.C3])"
    pi = sheet[1, 1].value
    ods.save()

    return HttpResponse("/gassdjango-api/gassdjangodataexport/data_export.py downloadODT() FUNCTIONALITY IS NOT YET COMPLETE :(\nHowever, the .ods was generated and is located @ /gassdjango-api/ODS_Export.ods", status=200)
    #return response

def downloadXLS():
    response = HttpResponse(content_type='application/vnd.ms-excel')
    response['Content-Disposition'] = 'attachment; filename="XLS_Export.xls"'

    wb = xlwt.Workbook(encoding='utf-8')
    ws = wb.add_sheet('gassdjangoproject')

    row = 0
    col = 0

    for i in range(0,85):
        ws.col(i).width = 30 * 256

    xls_data_array = []
    xls_data_array.append(["domingo","lunes","martes","miércoles","jueves","viernes","sábado"])
    xls_data_array.append(["sunday","monday","tuesday","wednesday","thursday","friday","saturday"])
    xls_data_array.append(["123","345","","","934"])
    xls_data_array.append(["","657","","","","092"])
    xls_data_array.append(["","","","234","934",""])

    for row, data in enumerate(xls_data_array):
        for attribute in data:
            ws.write(row, col, attribute, xlwt.easyxf("align: horiz left"))
            col += 1
        col = 0

    wb.save(response)
    return response

def downloadXLSX():
    response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
    response['Content-Disposition'] = 'attachment; filename="XLSX_Export.xlsx"'

    # Create an new Excel file and add a worksheet.
    workbook = xlsxwriter.Workbook(response)
    worksheet = workbook.add_worksheet('gassdjangoproject')

    # Widen the first column to make the text clearer.
    worksheet.set_column('A:A', 20)

    # Add a bold format to use to highlight cells.
    bold = workbook.add_format({'bold': True})

    # Write some simple text.
    worksheet.write('A1', 'Hello')

    # Text with formatting.
    worksheet.write('A2', 'World', bold)

    # Write some numbers, with row/column notation.
    worksheet.write(2, 0, 123)
    worksheet.write(3, 0, 123.456)

    # Insert an image.
    #worksheet.insert_image('B5', 'logo.png')

    workbook.close()

    return response

def downloadODP():
    response = HttpResponse(content_type="application/vnd.oasis.opendocument.presentation")
    response['Content-Disposition'] = 'attachment; filename="ODP_Export.odp"'

    odp = newdoc(doctype='odp', filename='ODP_Export.odp', template=None)
    odp.save()

    return HttpResponse("/gassdjango-api/gassdjangodataexport/data_export.py downloadODP() FUNCTIONALITY IS NOT YET COMPLETE :(\nHowever, the .odp was generated and is located @ /gassdjango-api/ODP_Export.odp", status=200)
    #return response

def downloadPPT():
    response = HttpResponse(content_type="application/vnd.ms-powerpoint")
    response['Content-Disposition'] = 'attachment; filename="PPT_Export.ppt"'
    prs = Presentation()
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Adding a Bullet Slide for gassdjangoproject'

    tf = body_shape.text_frame
    tf.text = 'Find the bullet slide layout'

    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.text for first bullet'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
    p.level = 2

    prs.save(response)
    return response

def downloadPPTX():
    response = HttpResponse(content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    response['Content-Disposition'] = 'attachment; filename="PPTX_Export.pptx"'
    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = 'Adding an AutoShape for gassdjangoproject'

    left = Inches(0.93)  # 0.93" centers this overall set of shapes
    top = Inches(3.0)
    width = Inches(1.75)
    height = Inches(1.0)

    shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
    shape.text = 'Step 1'

    left = left + width - Inches(0.4)
    width = Inches(2.0)  # chevrons need more width for visual balance

    for n in range(2, 6):
        shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
        shape.text = 'Step %d' % n
        left = left + width - Inches(0.4)

    prs.save(response)
    return response

def downloadPDF():
    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = 'attachment; filename="PDF_Export.pdf"'
    p = canvas.Canvas(response)
    pdf_text = "Here Is The Custom gassdjangoproject Text"
    p.drawString(100, 100, pdf_text)
    p.showPage()
    p.save()
    return response

def downloadMP3():
    response = HttpResponse(content_type='audio/mpeg')
    response['Content-Disposition'] = 'attachment; filename="MP3_Export.mp3"'
    mp3_text = "This example text for gassdjangoproject has been converted to an .MP3 file."
    mp3_data = gTTS(mp3_text, lang='en')
    mp3_data.write_to_fp(response)
    return response

def downloadHTM():
    htm_data = """
    <html>
    <head><title>gassdjangoproject HTM</title></head>
    <body>
    <h1>gassdjangoproject H1 heading</h1>
    <h3>gassdjangoproject H3 heading</h3>
    <h6>gassdjangoproject H6 heading</h6>
    <table border="1" cellspacing="10">
    <tr><td>Sight</td><td>Smell</td><td>Touch</td></tr>
    <tr><td>Taste</td><td>Sound</td><td>6th Sense</td></tr>
    </table>
    </body>
    </html>
    """
    response = HttpResponse(htm_data, content_type="text/html")
    response['Content-Disposition'] = 'attachment; filename="HTM_Export.htm"'
    return response

def downloadHTML():
    html_data = """
    <html>
    <head><title>gassdjangoproject HTML</title></head>
    <body>
    <h2>gassdjangoproject H2 heading</h1>
    <h4>gassdjangoproject H4 heading</h3>
    <h5>gassdjangoproject H5 heading</h6>
    <table border="1" cellspacing="10">
    <tr><td>Sight</td><td>Smell</td><td>Touch</td></tr>
    <tr><td>Taste</td><td>Sound</td><td>Sixth Sense</td></tr>
    </table>
    </body>
    </html>
    """
    response = HttpResponse(html_data, content_type="text/html")
    response['Content-Disposition'] = 'attachment; filename="HTML_Export.html"'
    return response

def downloadJSON():
    json_data = """[{"_id":{"":"000001"},"product_name":"Google Search Engine","company_name": "Google, Inc.","company_website":"https://www.google.com/"},{"_id":{"":"000002"},"product_name":"Yahoo! Search Engine","company_name":"Verizon Media","company_website":"https://www.yahoo.com/"},{"_id":{"":"000003"},"product_name":"Bing Search Engine","company_name":"Microsoft","company_website":"https://www.bing.com/"},{"_id":{"":"000004"},"product_name":"DuckDuckGo Search Engine","company_name":"Duck Duck Go, Inc.","company_website":"https://duckduckgo.com/"},{"_id":{"":"000005"},"product_name":"audible","company_name":"audible","company_website":"http://localhost:8000/"}]"""
    response = HttpResponse(json_data, content_type="application/json")
    response['Content-Disposition'] = 'attachment; filename="JSON_Export.json"'
    return response

def downloadSQL():
    sql_data = ""
    response = HttpResponse(sql_data, content_type="application/sql")
    response['Content-Disposition'] = 'attachment; filename="SQL_Export.sql"'
    return HttpResponse("/gassdjango-api/gassdjangodataexport/data_export.py downloadSQL() FUNCTIONALITY IS NOT YET COMPLETE :(", status=200)
    #return response

def downloadTXT():
    txt_data = "Example Text for gassdjangoproject .TXT File."
    response = HttpResponse(txt_data, content_type="text/plain")
    response['Content-Disposition'] = 'attachment; filename="TXT_Export.txt"'
    return response

def downloadRTF():
    rtf_data = """{Here is text for \b gassdjangoproject \b and more text \line\ The Farm \line\ \tBarn \line\ \t\tKittyKittyMeowMeow \line\ \t\tTweetyBird \line\ \tHouse \line\ \t\tDoggie \line\ \t\tPiggie \line\ Go to bed!}"""
    response = HttpResponse(rtf_data, content_type="application/rtf")
    response['Content-Disposition'] = 'attachment; filename="RTF_Export.rtf"'
    return response

def downloadJPG():
    response = HttpResponse(content_type='image/jpg')
    response['Content-Disposition'] = 'attachment; filename="JPG_Export.jpg"'
    img = Image.new('RGB',(500,500),color=(random.randint(0,255),random.randint(0,255),random.randint(0,255)))
    img_draw = ImageDraw.Draw(img)
    jpg_text = "JPG gassdjangoproject Text with get_random_word equal to %s" % (get_random_word(array_names).upper())
    img_draw.text((10,10),jpg_text,fill=(random.randint(0,255),random.randint(0,255),random.randint(0,255)))
    img.save(response, "JPEG")
    return response

def downloadJPEG():
    response = HttpResponse(content_type='image/jpeg')
    response['Content-Disposition'] = 'attachment; filename="JPEG_Export.jpeg"'
    img = Image.new('RGB',(500,500),color=(random.randint(0,255),random.randint(0,255),random.randint(0,255)))
    img_draw = ImageDraw.Draw(img)
    jpeg_text = "JPEG gassdjangoproject Text with get_random_word equal to %s" % (get_random_word(array_nouns).upper())
    img_draw.text((10,10),jpeg_text,fill=(random.randint(0,255),random.randint(0,255),random.randint(0,255)))
    img.save(response, "JPEG")
    return response

def downloadPNG():
    response = HttpResponse(content_type='image/png')
    response['Content-Disposition'] = 'attachment; filename="PNG_Export.png"'
    img = Image.new('RGB',(500,500),color=(random.randint(0,255),random.randint(0,255),random.randint(0,255)))
    img_draw = ImageDraw.Draw(img)
    png_text = "PNG gassdjangoproject Text with get_random_word equal to %s" % (get_random_word(array_verbs).upper())
    img_draw.text((10,10),png_text,fill=(random.randint(0,255),random.randint(0,255),random.randint(0,255)))
    img.save(response, "PNG")
    return response

def downloadGIF():
    response = HttpResponse(content_type='image/gif')
    response['Content-Disposition'] = 'attachment; filename="GIF_Export.gif"'

    images = []
    gif_text = "GIF gassdjangoproject Text"
    width = 300
    center = width // 2
    color_1 = (random.randint(0,255), random.randint(0,255), random.randint(0,255))
    color_2 = (random.randint(0,255), random.randint(0,255), random.randint(0,255))
    max_radius = int(center * 1.5)
    step = 8

    for i in range(0, max_radius, step):
        im = Image.new('RGB', (width, width), color_1)
        draw = ImageDraw.Draw(im)
        draw.ellipse((center - i, center - i, center + i, center + i), fill=color_2)
        images.append(im)
        draw.text((random.randint(10,150),random.randint(10,250)),gif_text,fill=(random.randint(0,255),random.randint(0,255),random.randint(0,255)))

    for i in range(0, max_radius, step):
        im = Image.new('RGB', (width, width), color_2)
        draw = ImageDraw.Draw(im)
        draw.ellipse((center - i, center - i, center + i, center + i), fill=color_1)
        images.append(im)
        draw.text((random.randint(10,150),random.randint(10,250)),gif_text,fill=(random.randint(0,255),random.randint(0,255),random.randint(0,255)))

    images[0].save(response, format='GIF', save_all=True, append_images=images[1:], optimize=False, duration=40, loop=0)

    return response

def downloadQRCode():
    img = qrcode.make('gassdjango')
    buf = BytesIO()
    img.save(buf)
    image_stream = buf.getvalue()
    response = HttpResponse(image_stream, content_type="image/png")
    return response
    
def downloadQRCodeUserInput(qrcode_text):
    img = qrcode.make(qrcode_text)
    buf = BytesIO()
    img.save(buf)
    image_stream = buf.getvalue()
    response = HttpResponse(image_stream, content_type="image/png")
    return response

def downloadZIP():
    b = BytesIO()
    z = ZipFile(b, 'w')
    z.write('ODS_Export.ods')
    z.write('ODT_Export.odt')
    z.write('ODP_Export.odp')
    z.close()
    b.seek(0)
    response = HttpResponse(b, content_type='application/zip')
    response['Content-Disposition'] = 'attachment; filename="ZIP_Export.zip"'
    return response
